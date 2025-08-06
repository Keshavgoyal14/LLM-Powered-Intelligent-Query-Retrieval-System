from urllib.parse import urlparse, parse_qs, unquote
import os
import tempfile
import requests
import zipfile
from langchain.schema import Document
from langchain_community.document_loaders import (
    PyPDFLoader, Docx2txtLoader, TextLoader, UnstructuredExcelLoader
)
from pptx import Presentation  # lightweight text extraction for pptx

def resolve_aspx(url: str) -> str:
    parsed = urlparse(url)
    if parsed.path.lower().endswith(".aspx"):
        qs = parse_qs(parsed.query)
        if "src" in qs:
            real_url = unquote(qs["src"][0])
            print(f"Resolved .aspx to: {real_url}")
            return real_url
    return url

def pptx_extract_text(pptx_path, url):
    """Extract text from pptx shapes (no OCR)."""
    text_content = []
    prs = Presentation(pptx_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_content.append(shape.text)
    combined_text = "\n".join(t.strip() for t in text_content if t.strip())
    if not combined_text:
        combined_text = "[No readable text found in pptx without OCR]"
    return [Document(page_content=combined_text, metadata={"source": url})]

def load_documents(url: str):
    print(f"Loading document from URL: {url}")
    url = resolve_aspx(url)

    parsed_url = urlparse(url)
    extension = os.path.splitext(parsed_url.path)[1].lower()

    allowed_ext = [
        ".pdf", ".docx", ".txt", ".pptx", ".xlsx", ".xls", ".png", ".jpg", ".jpeg"
    ]
    if extension not in allowed_ext:
        print(f"Unsupported file type: {extension}")
        return []

    max_size = 500 * 1024 * 1024

    try:
        head = requests.head(url, allow_redirects=True, timeout=10)
        content_length = head.headers.get("Content-Length")
        if content_length and int(content_length) > max_size:
            print(f"File too large: {content_length} bytes.")
            return []
    except Exception as e:
        print(f"HEAD request failed: {e}")
        content_length = None

    response = requests.get(url, timeout=60)
    response.raise_for_status()

    if not content_length and len(response.content) > max_size:
        print("File size exceeds limit after download")
        return []

    with tempfile.NamedTemporaryFile(delete=False, suffix=extension) as temp:
        temp.write(response.content)
        temp.flush()
        temp_path = temp.name

    try:
        if extension == ".pdf":
            docs = PyPDFLoader(temp_path).load()
        elif extension == ".docx":
            docs = Docx2txtLoader(temp_path).load()
        elif extension == ".txt":
            docs = TextLoader(temp_path).load()
        elif extension == ".pptx":
            docs = pptx_extract_text(temp_path, url)
        elif extension in [".xlsx", ".xls"]:
            docs = UnstructuredExcelLoader(temp_path).load()
        elif extension in [".png", ".jpg", ".jpeg"]:
            docs = [Document(page_content="[Image processing disabled]", metadata={"source": url})]
        else:
            docs = []
    finally:
        os.unlink(temp_path)

    return docs
